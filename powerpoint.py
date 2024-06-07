import image
import aspose.slides as slides
import image_format
from pptx import Presentation
import os
def PowerPoint(path):
    # open the file
    file_pptx = Presentation(path)
    imageIndex = 1
    all_text_powerpoint = ''
    # loop through slides
    for slide in file_pptx.slides:
        # loop through text in slides
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text_powerpoint += shape.text.lower().replace('\n', ' ')
   #=================================================================
    # load presentation
    with slides.Presentation(path) as pre:
        # loop through images
        for imag in pre.images:
            # format the path image
            file_name = "Image_{0}.{1}"
            #image_type = imag.content_type.split("/") # contain list [name image , type image]
            # contain type image
            image_type = imag.content_type.split("/")[1]
            # check if image_type in kinds of pictures
            image_formats = image_format.get_image_format(image_type)
            # save image
            imag.system_image.save(file_name.format(imageIndex, image_type), image_formats)
            # path the image
            path_of_image = file_name.format(imageIndex, image_type)
            # check if type image in image format
            if bool(image_format.ChiakImageFormat(path_of_image)):
                # read text from image
                img_text_png = image.IMAGES(path_of_image)
                # check if image contain on text
                if img_text_png != '':
                    # add text image in all text
                    all_text_powerpoint += ' ' + img_text_png
            # delete image
            os.remove(path_of_image)
            # number of image
            imageIndex += 1

    return all_text_powerpoint
